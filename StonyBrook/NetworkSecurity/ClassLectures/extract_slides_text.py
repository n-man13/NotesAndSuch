import sys
from pathlib import Path
from typing import List

folder = Path(__file__).parent

def extract_pdf(path: Path) -> str:
    try:
        from PyPDF2 import PdfReader
    except Exception as e:
        print(f"MISSING_PDF_LIB: {e}")
        return ""
    text_parts: List[str] = []
    try:
        reader = PdfReader(str(path))
        for p in reader.pages:
            try:
                t = p.extract_text()
            except Exception:
                t = None
            if t:
                text_parts.append(t)
    except Exception as e:
        text_parts.append(f"[ERROR reading PDF: {e}]")
    return "\n\n".join(text_parts)


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except Exception as e:
        print(f"MISSING_PPT_LIB: {e}")
        return ""
    text_parts: List[str] = []
    try:
        prs = Presentation(str(path))
        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_texts: List[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text:
                        slide_texts.append(shape.text)
                elif hasattr(shape, "text_frame") and shape.text_frame is not None:
                    slide_texts.append(shape.text_frame.text)
            if slide_texts:
                text_parts.append(f"--- Slide {slide_idx} ---\n" + "\n".join(slide_texts))
    except Exception as e:
        text_parts.append(f"[ERROR reading PPTX: {e}]")
    return "\n\n".join(text_parts)


def main():
    out_files = []
    for p in sorted(folder.iterdir()):
        if p.suffix.lower() == ".pdf":
            print(f"Extracting PDF: {p.name}")
            t = extract_pdf(p)
            if not t:
                t = "[no text extracted or missing library]"
            out = folder / (p.stem + ".txt")
            out.write_text(t, encoding="utf-8")
            out_files.append(out.name)
        elif p.suffix.lower() == ".pptx":
            print(f"Extracting PPTX: {p.name}")
            t = extract_pptx(p)
            if not t:
                t = "[no text extracted or missing library]"
            out = folder / (p.stem + ".txt")
            out.write_text(t, encoding="utf-8")
            out_files.append(out.name)
    print("WROTE:", out_files)

if __name__ == "__main__":
    main()
